from __future__ import annotations

from pathlib import Path

import pymupdf
from pptx import Presentation

from doc_to_abstract.config import Config, FileAnnotation
from doc_to_abstract.exceptions import FileExtractionError
from doc_to_abstract.template import read_template

PROMPTS_DIR = Path(__file__).resolve().parent.parent / "prompts"


def _format_annotation(annotations: dict[str, FileAnnotation], filepath: str) -> str:
    """Format annotation metadata for a file, if any."""
    ann = annotations.get(filepath) or annotations.get(Path(filepath).name)
    if not ann:
        return ""
    parts = []
    if ann.importance != "medium":
        parts.append(f"[Importance: {ann.importance}]")
    if ann.comment:
        parts.append(f"[Note: {ann.comment}]")
    return " ".join(parts)


def _extract_pdf(file_path: str, page_label: str) -> str:
    """Extract text from a PDF file using PyMuPDF."""
    try:
        doc = pymupdf.open(file_path)
    except Exception as e:
        raise FileExtractionError(f"Failed to open PDF: {file_path}: {e}") from e

    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            pages.append(f"--- {page_label} {i} ---\n{text}")
    doc.close()

    if not pages:
        raise FileExtractionError(f"No text could be extracted from: {file_path}")

    return "\n\n".join(pages)


def _extract_pptx(file_path: str, page_label: str) -> str:
    """Extract text from a PPTX file using python-pptx."""
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise FileExtractionError(f"Failed to open PPTX: {file_path}: {e}") from e

    pages = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            pages.append(f"--- {page_label} {i} ---\n" + "\n".join(texts))

    if not pages:
        raise FileExtractionError(f"No text could be extracted from: {file_path}")

    return "\n\n".join(pages)


def extract_text(file_path: str, page_label: str = "Page") -> str:
    """Extract text from a PDF or PPTX file."""
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx(file_path, page_label)
    return _extract_pdf(file_path, page_label)


def build_prompt(config: Config) -> str:
    """Build the full prompt for abstract generation."""
    template_path = PROMPTS_DIR / "system_prompt.txt"
    template = template_path.read_text(encoding="utf-8")

    # Format authors
    author_lines = []
    for a in config.authors:
        line = f"  - {a.name} ({a.affiliation})"
        if a.email:
            line += f" <{a.email}>"
        author_lines.append(line)
    authors_text = "\n".join(author_lines)

    # Build length constraint
    length_constraint = ""
    if config.max_words:
        length_constraint = f"- The abstract MUST NOT exceed {config.max_words} words."
    elif config.max_characters:
        length_constraint = f"- The abstract MUST NOT exceed {config.max_characters} characters."

    prompt = template.format(
        language=config.language,
        tone=config.tone,
        title=config.title,
        authors=authors_text,
        length_constraint=length_constraint,
    )

    # Append extra instructions
    if config.extra_instructions:
        extra_text = "\n".join(f"- {inst}" for inst in config.extra_instructions)
        prompt += f"\n\n## Additional Instructions\n{extra_text}\n"

    # Append extracted material text
    for i, mat_path in enumerate(config.materials, 1):
        mat_text = extract_text(mat_path)
        ann_text = _format_annotation(config.annotations, mat_path)
        if len(config.materials) > 1:
            header = f"\n\n## Main Material #{i} ({Path(mat_path).name})"
        else:
            header = f"\n\n## Main Material ({Path(mat_path).name})"
        if ann_text:
            header += f"\n{ann_text}"
        prompt += f"{header}\n{mat_text}\n"

    # Append reference papers
    for i, ref_path in enumerate(config.references, 1):
        ref_text = extract_text(ref_path)
        ann_text = _format_annotation(config.annotations, ref_path)
        header = f"\n## Reference Paper #{i} ({Path(ref_path).name})"
        if ann_text:
            header += f"\n{ann_text}"
        prompt += f"{header}\n{ref_text}\n"

    # Append supplementary materials
    for i, sup_path in enumerate(config.supplementary, 1):
        sup_text = extract_text(sup_path, page_label="Page")
        ann_text = _format_annotation(config.annotations, sup_path)
        header = f"\n## Supplementary Material #{i} ({Path(sup_path).name})"
        if ann_text:
            header += f"\n{ann_text}"
        prompt += f"{header}\n{sup_text}\n"

    # Append template content
    if config.template:
        template_text = read_template(config.template)
        prompt += (
            "\n## Submission Template\n"
            "The following is the submission template provided by the conference/workshop organizers. "
            "Use it to understand the required format, structure, length constraints, and any specific guidelines.\n\n"
            f"{template_text}\n"
        )

    return prompt
